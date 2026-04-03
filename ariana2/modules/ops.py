"""ops.py — file, system, image, screen, commands"""
import os, subprocess, zipfile, shutil, re, time, base64
from pathlib import Path
from io import BytesIO

# ── SCREEN ────────────────────────────────────────────────────
def capture_screen():
    try:
        import mss
        from PIL import Image
        with mss.mss() as s:
            mon = s.monitors[1]
            shot = s.grab(mon)
            img = Image.frombytes("RGB", shot.size, shot.bgra, "raw", "BGRX")
            img = img.resize((1280,720), Image.LANCZOS)
            buf = BytesIO()
            img.save(buf, format="PNG", optimize=True)
            return base64.b64encode(buf.getvalue()).decode()
    except Exception as e:
        return f"ERROR:{e}"

# ── FILES ─────────────────────────────────────────────────────
def list_dir(path="~"):
    path = os.path.expanduser(path)
    try:
        items = sorted(os.listdir(path))
        lines = []
        for it in items[:25]:
            full = os.path.join(path, it)
            ico = "📁" if os.path.isdir(full) else "📄"
            sz = ""
            if os.path.isfile(full):
                s = os.path.getsize(full)
                sz = f" ({s//1024}KB)" if s>1024 else f" ({s}B)"
            lines.append(f"{ico} {it}{sz}")
        return f"📂 {path}\n" + "\n".join(lines)
    except Exception as e:
        return f"❌ {e}"

def read_file(path):
    try:
        with open(os.path.expanduser(path),'r',encoding='utf-8',errors='replace') as f:
            return f.read(4000)
    except Exception as e:
        return f"❌ {e}"

def write_file(path, content):
    try:
        p = os.path.expanduser(path)
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        with open(p,'w',encoding='utf-8') as f:
            f.write(content)
        return f"✅ Saved: {p}"
    except Exception as e:
        return f"❌ {e}"

def delete_file(path):
    try:
        p = os.path.expanduser(path)
        shutil.rmtree(p) if os.path.isdir(p) else os.remove(p)
        return f"✅ Deleted: {p}"
    except Exception as e:
        return f"❌ {e}"

def compress(src, dst=None):
    p = Path(os.path.expanduser(src))
    if not p.exists(): return f"❌ Not found: {src}"
    dst = dst or str(p)+".zip"
    try:
        with zipfile.ZipFile(dst,'w',zipfile.ZIP_DEFLATED) as z:
            if p.is_file():
                z.write(p, p.name)
            else:
                for f in p.rglob('*'):
                    z.write(f, f.relative_to(p.parent))
        return f"✅ Zipped: {dst}"
    except Exception as e:
        return f"❌ {e}"

def open_file(path):
    try:
        subprocess.Popen(["xdg-open", os.path.expanduser(path)])
        return f"✅ Opening: {path}"
    except Exception as e:
        return f"❌ {e}"

# ── PDF ───────────────────────────────────────────────────────
def create_pdf(title, content, out=None):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import cm
        from reportlab.lib.colors import HexColor
        if not out:
            out = os.path.expanduser(f"~/Documents/ariana_{title[:15].replace(' ','_')}.pdf")
        os.makedirs(os.path.dirname(out), exist_ok=True)
        doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        ts = ParagraphStyle('T', parent=styles['Heading1'], fontSize=18, textColor=HexColor('#a855f7'), spaceAfter=16)
        bs = ParagraphStyle('B', parent=styles['Normal'], fontSize=11, leading=16, spaceAfter=6)
        story = [Paragraph(title, ts), Spacer(1,0.3*cm)]
        for line in content.split('\n'):
            if line.strip():
                safe = line.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
                story.append(Paragraph(safe, bs))
        doc.build(story)
        subprocess.Popen(["xdg-open", out])
        return f"✅ PDF ready: {out}"
    except ImportError:
        return "❌ pip install reportlab"
    except Exception as e:
        return f"❌ PDF error: {e}"

# ── PPT ───────────────────────────────────────────────────────
def create_ppt(title, slides, out=None):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        if not out:
            out = os.path.expanduser(f"~/Documents/ariana_{title[:15].replace(' ','_')}.pptx")
        prs = Presentation()
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = title
        for s in slides:
            sl2 = prs.slides.add_slide(prs.slide_layouts[1])
            sl2.shapes.title.text = s.get("title","")
            sl2.placeholders[1].text_frame.text = s.get("content","")
        prs.save(out)
        subprocess.Popen(["xdg-open", out])
        return f"✅ PPT ready: {out}"
    except ImportError:
        return "❌ pip install python-pptx"
    except Exception as e:
        return f"❌ PPT error: {e}"

# ── IMAGE GEN ─────────────────────────────────────────────────
def generate_image(prompt, out=None):
    import requests as req
    try:
        if not out:
            home = os.path.expanduser("~")
            os.makedirs(f"{home}/Pictures/ariana", exist_ok=True)
            out = f"{home}/Pictures/ariana/img_{int(time.time())}.jpg"
        url = f"https://image.pollinations.ai/prompt/{req.utils.quote(prompt)}?width=512&height=512&nologo=true"
        r = req.get(url, timeout=60)
        if r.status_code == 200:
            with open(out,"wb") as f: f.write(r.content)
            subprocess.Popen(["xdg-open", out])
            return f"✅ Image saved: {out}"
        return f"❌ HTTP {r.status_code}"
    except Exception as e:
        return f"❌ {e}"

# ── SYSTEM ────────────────────────────────────────────────────
def run_cmd(cmd):
    try:
        r = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=30)
        return (r.stdout or r.stderr or "✅ Done").strip()
    except subprocess.TimeoutExpired:
        return "❌ Timeout"
    except Exception as e:
        return f"❌ {e}"

def system_info():
    import psutil, platform
    cpu = psutil.cpu_percent(interval=0.5)
    mem = psutil.virtual_memory()
    disk = psutil.disk_usage('/')
    return (f"💻 {platform.system()} {platform.release()}\n"
            f"CPU: {cpu}%  RAM: {mem.percent}% ({mem.used//1024//1024}MB/{mem.total//1024//1024}MB)\n"
            f"Disk: {disk.percent}% used ({disk.free//1024//1024//1024}GB free)")

def get_ip():
    import socket, requests as req
    try:
        s=socket.socket(); s.connect(("8.8.8.8",80)); local=s.getsockname()[0]; s.close()
    except: local="N/A"
    try: pub = req.get("https://api.ipify.org",timeout=4).text
    except: pub="N/A"
    return f"🌐 Local: {local}  |  Public: {pub}"

# ── COMMAND ROUTER ─────────────────────────────────────────────
def route(text, ai):
    t = text.lower().strip()

    if any(x in t for x in ["screen dekh","screen kya","screen analyze","screen pe kya","screen batao"]):
        b64 = capture_screen()
        return ai.analyze_screen(b64, text) if not b64.startswith("ERROR") else b64

    if re.search(r"pdf bana|pdf banao", t):
        topic = re.sub(r"pdf bana(o)?","",text,flags=re.I).strip() or "Document"
        content = ai.chat(f'"{topic}" ke baare mein detailed notes likho.')
        return create_pdf(topic, content)

    if re.search(r"ppt bana|presentation bana|slides bana", t):
        topic = re.sub(r"(ppt|presentation|slides) bana(o)?","",text,flags=re.I).strip() or "Presentation"
        raw = ai.chat(f'"{topic}" pe 5 slides banao. Format exactly:\nSLIDE1: Title | Content\nSLIDE2: Title | Content\n...')
        slides = []
        for line in raw.split('\n'):
            if re.match(r'SLIDE\d+:',line,re.I):
                parts = line.split('|',1)
                slides.append({"title":re.sub(r'SLIDE\d+:\s*','',parts[0],flags=re.I).strip(),
                                "content":parts[1].strip() if len(parts)>1 else ""})
        return create_ppt(topic, slides or [{"title":topic,"content":raw[:200]}])

    if re.search(r"image bana|photo bana|image generate|image banao", t):
        prompt = re.sub(r"image bana(o)?|photo bana(o)?|image generate","",text,flags=re.I).strip() or "beautiful scenery"
        return generate_image(prompt)

    if re.search(r"files dikhao|folder dekho|list karo", t):
        path = (re.search(r'([~/][^\s]*)',text) or type('',(),{'group':lambda s,x:"~"})()).group(1)
        return list_dir(path)

    if re.search(r"system info|cpu|ram kitni|disk kitna", t):
        return system_info()

    if re.search(r"\bip\b.*\b(kya|mera|batao)\b|\b(mera|local|public) ip\b", t):
        return get_ip()

    if re.search(r"compress|zip karo", t):
        path = re.search(r'([~/][^\s]+)',text)
        return compress(path.group(1)) if path else "Kaunsi file compress karni hai?"

    if re.search(r"^(run |cmd:|terminal:|\$ )", t):
        cmd = re.sub(r"^(run |cmd:|terminal:|\$ )","",text,flags=re.I).strip()
        return run_cmd(cmd)

    if re.search(r"volume|awaaz kitni", t):
        lvl = re.search(r'\d+', t)
        return run_cmd(f"pactl set-sink-volume @DEFAULT_SINK@ {lvl.group()}%") if lvl else "Volume kitna karna hai? 0-100 bolo."

    if re.search(r"open |kholo ", t) and re.search(r'\w+\.\w+', t):
        url = re.sub(r"open |kholo ","",text,flags=re.I).strip()
        if not url.startswith("http"): url = "https://"+url
        subprocess.Popen(["xdg-open", url])
        return f"✅ Opening: {url}"

    if re.search(r"offline mode|ollama mode", t):
        ai.mode = "ollama"; return "💻 Offline mode — Ollama!"

    if re.search(r"online mode|groq mode", t):
        ai.mode = "groq"; return "🌐 Online mode — Groq!"

    if re.search(r"screenshot le", t):
        path = os.path.expanduser(f"~/Pictures/ss_{int(time.time())}.png")
        b64 = capture_screen()
        if not b64.startswith("ERROR"):
            import base64 as b6
            with open(path,"wb") as f: f.write(b6.b64decode(b64))
            return f"✅ Screenshot: {path}"
        return b64

    return None  # Let AI handle
